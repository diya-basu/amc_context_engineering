"""
document_extractors.py
=======================
Non-PDF extractors that mirror faiss_store.extract_pdf_text_full()'s output
shape, so build_index.py can treat every file type identically:

    [{"page_num": int, "text": str, "source": str, "extraction_method": str}, ...]

PDFs are NOT handled here — call faiss_store.extract_pdf_text_full() directly.
"""
from __future__ import annotations
import re
from pathlib import Path
from typing import Any, Dict, List


def extract_xlsx_text_full(path: str) -> List[Dict[str, Any]]:
    """Each sheet -> one 'page'. Rows rendered as markdown tables."""
    import openpyxl
    wb = openpyxl.load_workbook(path, data_only=True, read_only=True)
    pages = []
    for i, sheet_name in enumerate(wb.sheetnames, start=1):
        ws = wb[sheet_name]
        rows = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c).strip() if c is not None else "" for c in row]
            if any(cells):
                rows.append(cells)
        if not rows:
            continue
        header, body = rows[0], rows[1:]
        lines = ["| " + " | ".join(header) + " |",
                 "| " + " | ".join(["---"] * len(header)) + " |"]
        lines += ["| " + " | ".join(r) + " |" for r in body]
        text = f"[SHEET: {sheet_name}]\n" + "\n".join(lines)
        pages.append({
            "page_num": i, "text": text,
            "source": Path(path).name, "extraction_method": "openpyxl",
        })
    return pages


def extract_csv_text_full(path: str) -> List[Dict[str, Any]]:
    import csv
    with open(path, newline="", encoding="utf-8", errors="ignore") as f:
        rows = [r for r in csv.reader(f) if any(c.strip() for c in r)]
    if not rows:
        return []
    header, body = rows[0], rows[1:]
    lines = ["| " + " | ".join(header) + " |",
             "| " + " | ".join(["---"] * len(header)) + " |"]
    lines += ["| " + " | ".join(r) + " |" for r in body]
    return [{
        "page_num": 1, "text": "\n".join(lines),
        "source": Path(path).name, "extraction_method": "csv",
    }]


def extract_docx_text_full(path: str) -> List[Dict[str, Any]]:
    """One 'page' per ~40 paragraphs (docx has no real page concept)."""
    import docx
    d = docx.Document(path)
    paras = [p.text.strip() for p in d.paragraphs if p.text.strip()]

    tables_md = []
    for table in d.tables:
        rows = [[c.text.strip() for c in r.cells] for r in table.rows]
        if not rows:
            continue
        lines = ["| " + " | ".join(rows[0]) + " |",
                  "| " + " | ".join(["---"] * len(rows[0])) + " |"]
        lines += ["| " + " | ".join(r) + " |" for r in rows[1:]]
        tables_md.append("\n".join(lines))

    pages, chunk, PAGE_SIZE = [], [], 40
    for i in range(0, len(paras), PAGE_SIZE):
        chunk = paras[i:i + PAGE_SIZE]
        pages.append({
            "page_num": i // PAGE_SIZE + 1,
            "text": "\n\n".join(chunk),
            "source": Path(path).name, "extraction_method": "docx",
        })
    if tables_md:
        pages.append({
            "page_num": len(pages) + 1,
            "text": "\n\n".join(tables_md),
            "source": Path(path).name, "extraction_method": "docx_table",
        })
    return pages


def extract_pptx_text_full(path: str) -> List[Dict[str, Any]]:
    from pptx import Presentation
    prs = Presentation(path)
    pages = []
    for i, slide in enumerate(prs.slides, start=1):
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                parts.append(shape.text_frame.text.strip())
            if shape.has_table:
                rows = [[c.text.strip() for c in r.cells] for r in shape.table.rows]
                lines = ["| " + " | ".join(rows[0]) + " |",
                          "| " + " | ".join(["---"] * len(rows[0])) + " |"]
                lines += ["| " + " | ".join(r) + " |" for r in rows[1:]]
                parts.append("\n".join(lines))
        if shape_notes := getattr(slide, "notes_slide", None):
            if shape_notes.notes_text_frame.text.strip():
                parts.append(f"[NOTES] {shape_notes.notes_text_frame.text.strip()}")
        text = "\n\n".join(parts).strip()
        if text:
            pages.append({
                "page_num": i, "text": text,
                "source": Path(path).name, "extraction_method": "pptx",
            })
    return pages


def extract_txt_text_full(path: str) -> List[Dict[str, Any]]:
    text = Path(path).read_text(encoding="utf-8", errors="ignore")
    # split into ~1500 char "pages" on paragraph boundaries for parity with others
    paras = re.split(r"\n\s*\n", text)
    pages, buf, buf_len, pn = [], [], 0, 1
    for p in paras:
        p = p.strip()
        if not p:
            continue
        if buf_len + len(p) > 1500 and buf:
            pages.append({"page_num": pn, "text": "\n\n".join(buf),
                          "source": Path(path).name, "extraction_method": "txt"})
            pn += 1
            buf, buf_len = [], 0
        buf.append(p)
        buf_len += len(p)
    if buf:
        pages.append({"page_num": pn, "text": "\n\n".join(buf),
                      "source": Path(path).name, "extraction_method": "txt"})
    return pages


def extract_any(path: str, use_gemini: bool = True) -> List[Dict[str, Any]]:
    """Dispatcher — routes by extension. PDFs delegate to faiss_store."""
    ext = Path(path).suffix.lower()
    if ext == ".pdf":
        import faiss_store
        return faiss_store.extract_pdf_text_full(path, use_gemini=use_gemini)
    if ext in (".xlsx", ".xls"):
        return extract_xlsx_text_full(path)
    if ext == ".csv":
        return extract_csv_text_full(path)
    if ext == ".docx":
        return extract_docx_text_full(path)
    if ext == ".pptx":
        return extract_pptx_text_full(path)
    if ext == ".txt":
        return extract_txt_text_full(path)
    raise ValueError(f"Unsupported extension: {ext}")